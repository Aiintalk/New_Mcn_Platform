"""
app/services/document_parser.py

文档解析 service：支持 PDF / DOCX / XLSX / PPTX / TXT / MD 五种格式。

用法：
    from app.services.document_parser import parse_files_to_text
    raw_text = await parse_files_to_text(upload_files)

解析策略：
- 按扩展名分流到对应库
- 单文件失败：log warning + 跳过
- 全部失败或文本过短：抛 ValueError
- 合并文本截断到 _MAX_TEXT_LENGTH（8000 字符）
"""
import logging
import io
from typing import Any

from fastapi import UploadFile

logger = logging.getLogger(__name__)

_MAX_TEXT_LENGTH = 8000
_MIN_VALID_LENGTH = 10


async def parse_files_to_text(files: list[UploadFile]) -> str:
    """
    解析多个上传文件，合并为单一文本（截断 8000 字符）。

    Args:
        files: FastAPI UploadFile 列表

    Returns:
        合并后的文本，每文件以 `=== 文件: xxx ===\\n` 分隔

    Raises:
        ValueError: 无文件 / 全部解析失败 / 有效文本过短
    """
    items = await parse_files_to_items(files)
    texts = [f"=== 文件: {item['name']} ===\n{item['text']}" for item in items]

    combined = "\n\n".join(texts)

    if not combined.strip() or len(combined.strip()) < _MIN_VALID_LENGTH:
        raise ValueError(
            "无法从文件中提取有效文字内容，请尝试复制文档内容手动粘贴"
        )

    if len(combined) > _MAX_TEXT_LENGTH:
        combined = combined[:_MAX_TEXT_LENGTH]

    return combined


async def parse_files_to_items(files: list[UploadFile]) -> list[dict[str, str]]:
    """逐份解析上传文件，保留文件名与正文一一对应。"""
    if not files:
        raise ValueError("请上传文件")
    items: list[dict[str, str]] = []
    for f in files:
        filename = f.filename or "unknown"
        try:
            text = _extract_by_extension(filename, await f.read())
            if text and text.strip():
                items.append({"name": filename, "text": text})
            else:
                logger.warning("document_parser: file %s produced empty text, skipped", filename)
        except ValueError:
            raise
        except Exception as e:
            logger.warning("document_parser: failed to parse %s: %s", filename, e)
    combined_text = "\n".join(item["text"] for item in items).strip()
    if not combined_text or len(combined_text) < _MIN_VALID_LENGTH:
        raise ValueError("无法从文件中提取有效文字内容，请尝试复制文档内容手动粘贴")
    return items


def parse_file_content_to_item(filename: str, content: bytes) -> dict[str, str]:
    """解析已读取的单份文档，供需要先限制上传大小的路由复用。"""
    text = _extract_by_extension(filename, content)
    if not text or not text.strip():
        raise ValueError("无法从文件中提取有效文字内容，请尝试复制文档内容手动粘贴")
    return {"name": filename, "text": text}


def _extract_by_extension(filename: str, content_bytes: bytes) -> str:
    """按文件扩展名分流到对应解析器。"""
    lower = filename.lower()

    if lower.endswith(".pdf"):
        return _extract_pdf(content_bytes)
    if lower.endswith(".docx"):
        return _extract_docx(content_bytes)
    if lower.endswith(".xlsx"):
        return _extract_xlsx(content_bytes)
    if lower.endswith(".xls"):
        raise ValueError(
            "不支持 .xls 老格式，请转换为 .xlsx 后上传"
        )
    if lower.endswith(".pptx"):
        return _extract_pptx(content_bytes)
    # .txt, .md, .csv, and unknown extensions: try utf-8 decode
    return _extract_text(content_bytes)


def _extract_pdf(content_bytes: bytes) -> str:
    """用 pypdf 提取 PDF 文本。"""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content_bytes))
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        pages.append(text)
    return "\n".join(pages)


def _extract_docx(content_bytes: bytes) -> str:
    """用 python-docx 提取 DOCX 文本。"""
    from docx import Document

    doc = Document(io.BytesIO(content_bytes))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(content_bytes: bytes) -> str:
    """用 openpyxl 提取 XLSX 文本（所有 sheet 转 CSV 拼接）。"""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(",".join(cells))
        sheets.append(f"[{sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets)


def _extract_pptx(content_bytes: bytes) -> str:
    """用 python-pptx 提取 PPTX 文本（所有 slide 的 text frame 拼接）。"""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content_bytes))
    slides_text = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides_text.append(" ".join(texts))
    return "\n".join(slides_text)


def _extract_text(content_bytes: bytes | str) -> str:
    """UTF-8 解码文本文件。"""
    if isinstance(content_bytes, str):
        return content_bytes
    try:
        return content_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return content_bytes.decode("gbk", errors="replace")
